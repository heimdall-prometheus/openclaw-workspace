#!/usr/bin/env python3
"""
Becker Sicherheitstechnik - Odoo ERP Projektplan
PowerPoint basierend auf Aufwandschätzung_v1.xlsx
"""

import sys
sys.path.insert(0, '/home/reisig/.openclaw/workspace/.venv/lib/python3.12/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Farben
PRIMARY = RGBColor(0x1a, 0x56, 0xdb)
SECONDARY = RGBColor(0x10, 0xb9, 0x81)
ACCENT = RGBColor(0xf5, 0x9e, 0x0b)
DARK = RGBColor(0x1f, 0x29, 0x37)
LIGHT = RGBColor(0xf3, 0xf4, 0xf6)
WHITE = RGBColor(0xff, 0xff, 0xff)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(prs, title, emoji=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    
    if emoji:
        emoji_box = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(2), Inches(1))
        tf = emoji_box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(72)
        p.alignment = PP_ALIGN.CENTER
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, items, highlight_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = highlight_color or PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)
    return slide

def add_table_slide(prs, title, headers, rows, highlight_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = highlight_color or PRIMARY
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    cols = len(headers)
    table_rows = len(rows) + 1
    row_height = min(0.45, 4.5 / table_rows)
    table = slide.shapes.add_table(table_rows, cols, Inches(0.2), Inches(1.4), Inches(9.6), Inches(row_height * table_rows)).table
    
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.font.size = Pt(11)
    
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.color.rgb = DARK
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
    return slide

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Titel
    add_title_slide(prs, 
        "Becker Sicherheitstechnik", 
        "Odoo ERP Projektplan\nBasierend auf Aufwandschätzung v1")
    
    # Strategie
    add_content_slide(prs, "Strategischer Ansatz", [
        "Quick Wins zuerst → Gute Stimmung im Unternehmen",
        "Projekt im guten Licht erscheinen lassen",
        "Komplexe Dinge nach hinten schieben",
        "Frühe Erfolge schaffen Buy-In für schwierige Phasen",
        "Gesamtaufwand: 63-95 Personentage (~79 PT Mittel)",
    ])
    
    # Gesamtübersicht
    add_table_slide(prs, "Aufwandsübersicht nach Phasen",
        ["Phase", "Min PT", "Max PT", "Mittel"],
        [
            ["1. Analyse & Projektsetup", "9", "14", "11,5"],
            ["2. Infrastruktur & Basis", "3", "5", "4"],
            ["3. Stammdatenmigration", "2", "3", "2,5"],
            ["4. Buchhaltung", "8,5", "12,5", "10,5"],
            ["5. Einkauf & Beschaffung", "4", "7", "5,5"],
            ["6. Lager & Logistik", "5", "9", "7"],
            ["7. Vertrieb & Auftragsabwicklung", "4", "5,5", "4,75"],
            ["8. Projektmanagement & CRM", "6", "8", "7"],
            ["9. Mängeldoku & Service", "7", "9", "8"],
            ["10. Custom Module", "7", "11", "9"],
            ["11. Schulung & Change", "5,5", "8", "6,75"],
            ["12. Test & Rollout", "2", "3", "2,5"],
            ["GESAMT", "63", "95", "~79"],
        ])
    
    # Phase 1
    add_section_slide(prs, "Phase 1: Analyse & Setup", "📋")
    add_table_slide(prs, "Analyse & Projektsetup (Woche 1-3)",
        ["Aufgabe", "Min PT", "Max PT", "Deliverable"],
        [
            ["Kick-off & Projektorganisation", "2", "3", "Projektteam steht"],
            ["Prozessworkshops (Ist/Soll)", "4", "6", "Prozesse dokumentiert"],
            ["Anforderungs-Mapping & Backlog", "2", "3", "Priorisiertes Backlog"],
            ["Projektplan & Roadmap", "1", "2", "Verbindlicher Zeitplan"],
        ], PRIMARY)
    
    # Phase 2+3
    add_table_slide(prs, "Infrastruktur & Stammdaten (Woche 4-6)",
        ["Aufgabe", "Min PT", "Max PT", "Deliverable"],
        [
            ["Installation Odoo Community", "0,5", "0,5", "System läuft"],
            ["Server-Konfiguration", "0,5", "0,5", "Produktivumgebung"],
            ["Integration OCA-Bibliotheken", "1", "2", "Module installiert"],
            ["Rechte- & Rollenkonzept", "1", "2", "User angelegt"],
            ["Stammdatenmigration", "2", "3", "Echte Daten im System"],
        ], SECONDARY)
    
    # Quick Wins
    add_section_slide(prs, "Quick Wins: Vertrieb & Einkauf", "⚡")
    
    add_table_slide(prs, "Vertrieb & Auftragsabwicklung (Woche 6-8)",
        ["Aufgabe", "Min PT", "Max PT", "Deliverable"],
        [
            ["Angebots- & Auftragsprozesse", "1", "2", "Erstes Angebot aus Odoo"],
            ["Kundenspezifische Preisregeln", "0,5", "0,5", "Individuelle Preise"],
            ["Auftragsbestätigung & Kopierlogik", "0,5", "1", "Effiziente Aufträge"],
            ["POS-Modul", "2", "2", "Thekenverkauf funktioniert"],
        ], SECONDARY)
    
    add_table_slide(prs, "Einkauf & Beschaffung (Woche 8-10)",
        ["Aufgabe", "Min PT", "Max PT", "Deliverable"],
        [
            ["Bestellung & Wareneingang", "1", "2", "Bestellungen digital"],
            ["Lieferantenverwaltung", "0,5", "0,5", "Lieferanten gepflegt"],
            ["Übersicht offene Bestellungen", "0,5", "0,5", "Transparenz"],
            ["3-Wege-Prüfung", "2", "4", "Automatischer Abgleich"],
        ], SECONDARY)
    
    # Lager & Buchhaltung
    add_section_slide(prs, "Kernprozesse: Lager & Finanzen", "🔧")
    
    add_table_slide(prs, "Lager & Logistik (Woche 10-12)",
        ["Aufgabe", "Min PT", "Max PT", "Zusatz"],
        [
            ["Mehrere Lagerorte", "0,5", "2", "Eigenleistung möglich"],
            ["Bestandsführung & Umlagerung", "0,5", "1", ""],
            ["Scanner-Integration", "2", "3", "Bei mobilen Scannern"],
            ["Etikettendruck", "1", "2", "Modul ~500€"],
            ["Meldebestand", "1", "1", ""],
        ], PRIMARY)
    
    add_table_slide(prs, "Buchhaltung (Woche 12-15)",
        ["Aufgabe", "Min PT", "Max PT", "Zusatz"],
        [
            ["DATEV-Schnittstelle", "1", "2", "Modul ~500€"],
            ["X-Rechnung / ZUGFeRD", "0,5", "0,5", "inkl. OCR"],
            ["Eingangsrechnungsprüfung", "3", "4", "Automatisiert"],
            ["Pflichtfelder & Validierung", "2", "3", ""],
            ["Kontenplan & Migration", "2", "3", ""],
        ], PRIMARY)
    
    # Projektmanagement
    add_table_slide(prs, "Projektmanagement & CRM (Woche 15-17)",
        ["Aufgabe", "Min PT", "Max PT", "Deliverable"],
        [
            ["Projektanlage & Nachkalkulation", "2", "3", "Kosten sichtbar"],
            ["Projektdokumentation", "1", "2", "Alles am Projekt"],
            ["CRM & Kundenhistorie", "2", "2", "Komplette Kundenakte"],
            ["Profitabilität je Projekt", "1", "1", "Gewinn auf Knopfdruck"],
        ], ACCENT)
    
    # Komplexe Phasen
    add_section_slide(prs, "Komplexe Module", "🔗")
    
    add_table_slide(prs, "Mängeldoku & Service (Woche 17-20)",
        ["Aufgabe", "Min PT", "Max PT", "Deliverable"],
        [
            ["Wartungsverträge", "1", "1", "Verträge digital"],
            ["Mobile App Monteure", "4", "5", "Monteure arbeiten mobil"],
            ["Mängeldoku mit Foto", "1", "2", "Fotos am Projekt"],
            ["Digitale Unterschrift", "1", "1", "Kunde signiert digital"],
        ], DARK)
    
    add_table_slide(prs, "Custom Module & Schnittstellen (Woche 20-23)",
        ["Aufgabe", "Min PT", "Max PT", "Deliverable"],
        [
            ["S-Firm Schnittstelle", "1", "2", "Überweisungen automatisiert"],
            ["GAEB-Import", "3", "4", "Ausschreibungen → Angebote"],
            ["GAEB Mapping (N:1)", "-", "-", "Statistik & Abrechnung"],
            ["Sonderlösungen", "3", "5", "Custom Workflows"],
        ], DARK)
    
    # Schulung & Go-Live
    add_section_slide(prs, "Schulung & Go-Live", "🚀")
    
    add_table_slide(prs, "Schulung & Change (Woche 23-25)",
        ["Aufgabe", "Min PT", "Max PT", "Deliverable"],
        [
            ["Key-User-Training", "3", "4", "Key User fit"],
            ["Admin-Training", "0,5", "1", "Admins können verwalten"],
            ["Endanwender-Workshops", "1", "2", "Alle geschult"],
            ["Schulungsunterlagen", "1", "1", "Dokumentation fertig"],
        ], SECONDARY)
    
    add_table_slide(prs, "Test & Rollout (Woche 25-26)",
        ["Aufgabe", "Min PT", "Max PT", "Deliverable"],
        [
            ["Integrationstests", "-", "-", "Prozesse getestet"],
            ["Pilotbetrieb", "-", "-", "Echtbetrieb mit Fallback"],
            ["Bugfixes & Nachjustierung", "2", "3", "System stabil"],
            ["Go-Live-Begleitung", "-", "-", "Support beim Start"],
        ], SECONDARY)
    
    # Meilensteine
    add_table_slide(prs, "Meilensteine & Timeline",
        ["Woche", "Meilenstein", "Ergebnis"],
        [
            ["3", "Analyse abgeschlossen", "Anforderungen klar"],
            ["5", "System steht", "Odoo läuft mit Daten"],
            ["8", "Vertrieb produktiv", "Angebote aus Odoo"],
            ["10", "Einkauf produktiv", "Bestellungen digital"],
            ["15", "Buchhaltung komplett", "X-Rechnung ready"],
            ["17", "Projektcontrolling", "Kein Umsatzverlust"],
            ["20", "Monteure digital", "Mobile App live"],
            ["23", "Custom fertig", "S-Firm & GAEB"],
            ["25", "Schulungen done", "Team bereit"],
            ["26", "GO-LIVE", "Produktiver Betrieb"],
        ])
    
    # Kosten
    add_content_slide(prs, "Zusatzkosten (Module)", [
        "DATEV-Schnittstelle: ~500€",
        "Etikettendruck-Modul: ~500€",
        "Summe Module: ~1.000€",
        "",
        "Scanner-Hardware: Nach Bedarf (separate Beschaffung)",
    ])
    
    # Risiken
    add_table_slide(prs, "Risiken & Mitigation",
        ["Risiko", "Wahrscheinlichkeit", "Mitigation"],
        [
            ["Datenqualität Migration", "Mittel", "Frühe Testmigration"],
            ["GAEB-Komplexität", "Hoch", "Puffer, iterativ"],
            ["Akzeptanz Mitarbeiter", "Niedrig", "Quick Wins → Buy-In"],
            ["Scanner-Hardware", "Mittel", "Frühzeitig beschaffen"],
        ])
    
    # Nächste Schritte
    add_content_slide(prs, "Nächste Schritte", [
        "1. Kick-off Meeting terminieren",
        "2. Key User pro Bereich benennen",
        "3. Server-Infrastruktur vorbereiten",
        "4. Datenexport aus Altsystem starten",
        "5. Hardware-Anforderungen klären (Scanner, Drucker)",
    ])
    
    # Abschluss
    add_title_slide(prs, "Fragen?", "Becker Sicherheitstechnik\nOdoo ERP Projekt 2026")
    
    output_path = "/home/reisig/.openclaw/workspace/projects/becker-odoo/Projektplan_Becker_Odoo_Final.pptx"
    prs.save(output_path)
    print(f"✅ Präsentation erstellt: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
