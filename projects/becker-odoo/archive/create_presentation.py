#!/usr/bin/env python3
"""
Becker Sicherheitstechnik - Odoo ERP Projektplanung
PowerPoint Präsentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from datetime import datetime

# Farben
PRIMARY = RGBColor(0x1a, 0x56, 0xdb)  # Blau
SECONDARY = RGBColor(0x10, 0xb9, 0x81)  # Grün
ACCENT = RGBColor(0xf5, 0x9e, 0x0b)  # Orange
DARK = RGBColor(0x1f, 0x29, 0x37)  # Dunkelgrau
LIGHT = RGBColor(0xf3, 0xf4, 0xf6)  # Hellgrau
WHITE = RGBColor(0xff, 0xff, 0xff)

def add_title_slide(prs, title, subtitle):
    """Titelfolie"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Hintergrund-Shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    
    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Untertitel
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title, emoji=""):
    """Abschnitts-Trennfolie"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Hintergrund
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    
    # Emoji
    if emoji:
        emoji_box = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(2), Inches(1))
        tf = emoji_box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(72)
        p.alignment = PP_ALIGN.CENTER
    
    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, items, highlight_color=None):
    """Standard-Inhaltsfolie"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titel-Balken
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = highlight_color or PRIMARY
    header.line.fill.background()
    
    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK
        p.space_after = Pt(12)
    
    return slide

def add_table_slide(prs, title, headers, rows, highlight_color=None):
    """Folie mit Tabelle"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titel-Balken
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = highlight_color or PRIMARY
    header_shape.line.fill.background()
    
    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Tabelle
    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.3), Inches(1.5), Inches(9.4), Inches(0.5 * table_rows)).table
    
    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.font.size = Pt(14)
    
    # Rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(13)
            para.font.color.rgb = DARK
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
    
    return slide

def add_milestone_slide(prs):
    """Meilensteine visuell"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titel-Balken
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Meilensteine & Timeline"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Timeline Linie
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(9), Inches(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT
    line.line.fill.background()
    
    milestones = [
        ("Woche 4", "Go-Live Basis", "⚡", SECONDARY),
        ("Woche 6", "Projektcontrolling", "📊", PRIMARY),
        ("Woche 10", "Kernprozesse", "🔧", PRIMARY),
        ("Woche 14", "Optimierung", "🚀", ACCENT),
        ("Woche 20", "Vollausbau", "🔗", DARK),
    ]
    
    for i, (week, desc, emoji, color) in enumerate(milestones):
        x = Inches(0.5 + i * 1.9)
        
        # Punkt
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), Inches(2.85), Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        
        # Woche
        week_box = slide.shapes.add_textbox(x, Inches(2.3), Inches(1.5), Inches(0.5))
        tf = week_box.text_frame
        p = tf.paragraphs[0]
        p.text = week
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        # Beschreibung
        desc_box = slide.shapes.add_textbox(x, Inches(3.3), Inches(1.5), Inches(1))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{emoji}\n{desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Hauptfunktion"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Titel
    add_title_slide(prs, 
        "Becker Sicherheitstechnik", 
        "Odoo ERP Projektplanung\n02.02.2026")
    
    # Strategie
    add_content_slide(prs, "Strategischer Ansatz", [
        "Quick Wins zuerst → Gute Stimmung im Unternehmen",
        "Projekt im guten Licht erscheinen lassen",
        "Komplexe Dinge nach hinten schieben",
        "Frühe Erfolge schaffen Buy-In für schwierige Phasen",
        "Change Management: Momentum aufbauen!",
    ])
    
    # Phasen-Übersicht
    add_table_slide(prs, "4 Phasen zum Erfolg",
        ["Phase", "Fokus", "Zeitraum", "Charakter"],
        [
            ["Phase 1", "Quick Wins", "Wochen 1-4", "⚡ Schnelle Erfolge"],
            ["Phase 2", "Kernprozesse", "Wochen 5-10", "🔧 Solide Basis"],
            ["Phase 3", "Optimierung", "Wochen 11-14", "🚀 Effizienz"],
            ["Phase 4", "Custom & Integration", "Wochen 15-20", "🔗 Speziallösungen"],
        ])
    
    # Phase 1
    add_section_slide(prs, "PHASE 1: QUICK WINS", "⚡")
    
    add_table_slide(prs, "AP 1.1: Stammdaten & Grundstruktur (Woche 1-2)",
        ["Task", "Ergebnis"],
        [
            ["Odoo Instanz aufsetzen", "✅ System läuft"],
            ["Benutzer anlegen + Rollen", "✅ Alle können einloggen"],
            ["3 Standorte konfigurieren", "✅ Nummernkreise getrennt"],
            ["Artikelstamm importieren", "✅ Top 100 Produkte verfügbar"],
            ["Kundenstamm importieren", "✅ Top 50 Kunden drin"],
        ], SECONDARY)
    
    add_table_slide(prs, "AP 1.2-1.4: Angebote, X-Rechnung, Lager",
        ["Arbeitspaket", "Woche", "Deliverable"],
        [
            ["Angebote & Aufträge", "2-3", "Erstes echtes Angebot aus Odoo"],
            ["X-Rechnung (Gesetzlich!)", "3-4", "Gesetzeskonforme E-Rechnung"],
            ["Basis-Lagerhaltung", "4", "Echtzeit-Bestandsübersicht"],
        ], SECONDARY)
    
    add_content_slide(prs, "Phase 1 Erfolge (nach 4 Wochen)", [
        "✅ System steht mit echten Daten",
        "✅ Angebote/Aufträge laufen digital",
        "✅ X-Rechnung → Gesetzeskonform",
        "✅ Lagerbestand → Transparenz",
        "✅ Team hat erste Erfolge → MOTIVATION!",
    ], SECONDARY)
    
    # Phase 2
    add_section_slide(prs, "PHASE 2: KERNPROZESSE", "🔧")
    
    add_table_slide(prs, "AP 2.1: Projekt-Modul (KRITISCH!)",
        ["Task", "Ergebnis"],
        [
            ["Projektstruktur definieren", "✅ Einheitliche Nomenklatur"],
            ["Projekte anlegen (5 Piloten)", "✅ Echte Projekte live"],
            ["Kosten auf Projekt buchen", "✅ Material + Zeit zugeordnet"],
            ["Weiterberechnungs-Workflow", "✅ KEIN UMSATZVERLUST MEHR!"],
        ], PRIMARY)
    
    add_table_slide(prs, "AP 2.2-2.5: Einkauf, Zeit, POS, DATEV",
        ["Arbeitspaket", "Woche", "Deliverable"],
        [
            ["Einkauf & 3-Wege-Prüfung", "6-7", "Automatische Rechnungsprüfung"],
            ["Zeiterfassung", "7-8", "Zettelwirtschaft beendet"],
            ["POS / Thekenverkauf", "8-9", "Digitale Kasse"],
            ["DATEV Export", "9-10", "Automatischer Buchungsexport"],
        ], PRIMARY)
    
    # Phase 3
    add_section_slide(prs, "PHASE 3: OPTIMIERUNG", "🚀")
    
    add_table_slide(prs, "Optimierung (Wochen 11-14)",
        ["Arbeitspaket", "Woche", "Deliverable"],
        [
            ["Monteur-Workflow (Basis)", "11-12", "Strukturierte Planung"],
            ["Kundenpreise & Konditionen", "12-13", "Historie mit Preisinfos"],
            ["Wartungsverträge (Basis)", "13-14", "Digital verwaltet"],
        ], ACCENT)
    
    # Phase 4
    add_section_slide(prs, "PHASE 4: CUSTOM & INTEGRATION", "🔗")
    
    add_table_slide(prs, "Custom Development (Wochen 15-20)",
        ["Arbeitspaket", "Woche", "Deliverable"],
        [
            ["S-Firm Schnittstelle", "15-16", "Automatisierte Überweisungen"],
            ["GAEB Import", "16-18", "Ausschreibungen → Angebote"],
            ["Mobile App Monteure", "18-20", "Monteure papierlos"],
        ], DARK)
    
    # Timeline
    add_milestone_slide(prs)
    
    # Aufwand
    add_table_slide(prs, "Aufwandsübersicht",
        ["Phase", "Wochen", "Tage Aufwand", "Charakter"],
        [
            ["Phase 1", "1-4", "9 Tage", "Quick Wins ⚡"],
            ["Phase 2", "5-10", "14 Tage", "Kernprozesse 🔧"],
            ["Phase 3", "11-14", "9 Tage", "Optimierung 🚀"],
            ["Phase 4", "15-20", "14 Tage", "Custom 🔗"],
            ["GESAMT", "20", "~46 Tage", ""],
        ])
    
    # Risiken
    add_table_slide(prs, "Risiken & Mitigation",
        ["Risiko", "Wahrscheinlichkeit", "Mitigation"],
        [
            ["Datenqualität Import", "Mittel", "Früh testen, parallel bereinigen"],
            ["Widerstand Mitarbeiter", "Niedrig", "Quick Wins → Buy-In"],
            ["GAEB-Komplexität", "Hoch", "Puffer eingeplant, iterativ"],
            ["S-Firm Format", "Mittel", "Frühe Klärung mit Bank"],
        ])
    
    # Nächste Schritte
    add_content_slide(prs, "Nächste Schritte", [
        "1. Workshop Kickoff → Prozesse validieren",
        "2. Key User benennen → Pro Bereich Ansprechpartner",
        "3. Instanz aufsetzen → Technische Vorbereitung",
        "4. Datenexport Altsystem → Für Import vorbereiten",
        "5. Termine fixieren → Kalender abstimmen",
    ])
    
    # Abschluss
    add_title_slide(prs, 
        "Fragen?", 
        "Erik Reisig | Becker Sicherheitstechnik\n02.02.2026")
    
    # Speichern
    output_path = "/home/reisig/.openclaw/workspace/projects/becker-odoo/Projektplanung_Becker_Odoo.pptx"
    prs.save(output_path)
    print(f"✅ Präsentation erstellt: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
