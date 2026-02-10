#!/usr/bin/env python3
"""
Becker Sicherheitstechnik - Odoo ERP Projektplan v3
NEUE REIHENFOLGE: PM/CRM direkt nach Einkauf, Lager/Buchhaltung danach
"""

import sys
sys.path.insert(0, '/home/reisig/.openclaw/workspace/.venv/lib/python3.12/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Farben
PRIMARY = RGBColor(0x1a, 0x56, 0xdb)
SECONDARY = RGBColor(0x10, 0xb9, 0x81)
ACCENT = RGBColor(0xf5, 0x9e, 0x0b)
DARK = RGBColor(0x1f, 0x29, 0x37)
LIGHT = RGBColor(0xf3, 0xf4, 0xf6)
WHITE = RGBColor(0xff, 0xff, 0xff)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(prs, title, emoji=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    
    if emoji:
        emoji_box = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(2), Inches(1))
        tf = emoji_box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(72)
        p.alignment = PP_ALIGN.CENTER
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, items, highlight_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = highlight_color or PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)
    return slide

def add_table_slide(prs, title, headers, rows, highlight_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = highlight_color or PRIMARY
    header_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    cols = len(headers)
    table_rows = len(rows) + 1
    row_height = min(0.45, 4.5 / table_rows)
    table = slide.shapes.add_table(table_rows, cols, Inches(0.2), Inches(1.4), Inches(9.6), Inches(row_height * table_rows)).table
    
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.font.size = Pt(11)
    
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.color.rgb = DARK
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
    return slide

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Titel
    add_title_slide(prs, 
        "Becker Sicherheitstechnik", 
        "Odoo ERP Projektplan\n9. Februar - 20. Juli 2026")
    
    # Prozess-Reihenfolge (NEUE VERSION)
    add_table_slide(prs, "Prozess-Reihenfolge",
        ["#", "Prozessbereich", "Strategie"],
        [
            ["1", "Vertrieb & Angebote", "Quick Win - sofort sichtbarer Nutzen"],
            ["2", "Einkauf & Bestellung", "Gemeinsame Stammdaten nutzen"],
            ["3", "Projektmanagement", "Aufträge → Projekte verknüpfen"],
            ["4", "CRM & Kundenakte", "Komplette Kundenhistorie"],
            ["5", "Lager & Bestand", "Basis für Warenfluss"],
            ["6", "Buchhaltung & DATEV", "Finanzen parallel aufbauen"],
            ["7", "Monteur-App & Service", "Mobile Lösung Außendienst"],
            ["8", "Schnittstellen", "S-Firm, GAEB (komplex)"],
        ])
    
    # Timeline
    add_table_slide(prs, "Timeline: 23 Wochen @ 3.5 PT/Woche",
        ["Wochen", "Zeitraum", "Phase", "PT"],
        [
            ["1-2", "9.-21. Feb", "Analyse + Basis-Setup", "8"],
            ["3-5", "22. Feb - 14. Mär", "Vertrieb & Angebote", "12"],
            ["6-8", "15. Mär - 4. Apr", "Einkauf & Bestellung", "8"],
            ["9-11", "5.-25. Apr", "Projektmanagement", "10"],
            ["12-14", "26. Apr - 16. Mai", "CRM & Kundenakte", "8"],
            ["15-17", "17. Mai - 6. Jun", "Lager & Bestand", "10"],
            ["18-19", "7.-20. Jun", "Buchhaltung & DATEV", "12"],
            ["20-21", "21. Jun - 4. Jul", "Monteur-App & Service", "8"],
            ["22-23", "5.-20. Jul", "Schnittstellen & Go-Live", "12"],
        ])
    
    # Quick Wins Phase
    add_section_slide(prs, "Phase 1-2: Quick Wins", "⚡")
    
    add_table_slide(prs, "Vertrieb & Angebote (Woche 3-5)",
        ["Aufgabe", "PT", "Deliverable"],
        [
            ["Angebots- & Auftragsprozesse", "4", "Erstes Angebot aus Odoo"],
            ["Kundenspezifische Preisregeln", "2", "Individuelle Preise"],
            ["Auftragsbestätigung & Kopierlogik", "2", "Effiziente Aufträge"],
            ["POS-Modul (Thekenverkauf)", "4", "Direktverkauf funktioniert"],
        ], SECONDARY)
    
    add_table_slide(prs, "Einkauf & Bestellung (Woche 6-8)",
        ["Aufgabe", "PT", "Deliverable"],
        [
            ["Bestellung & Wareneingang", "2", "Bestellungen digital"],
            ["Lieferantenverwaltung", "1", "Lieferanten gepflegt"],
            ["Übersicht offene Bestellungen", "1", "Transparenz"],
            ["3-Wege-Prüfung", "4", "Automatischer Abgleich"],
        ], SECONDARY)
    
    # Projekte & CRM (FRÜHER)
    add_section_slide(prs, "Phase 3-4: Projekte & CRM", "📊")
    
    add_table_slide(prs, "Projektmanagement (Woche 9-11)",
        ["Aufgabe", "PT", "Deliverable"],
        [
            ["Projektanlage & Nachkalkulation", "3", "Kosten sichtbar"],
            ["Projektdokumentation", "2", "Alles am Projekt"],
            ["Profitabilität je Projekt", "3", "Gewinn auf Knopfdruck"],
            ["Reports & Dashboards", "2", "Management-Überblick"],
        ], ACCENT)
    
    add_table_slide(prs, "CRM & Kundenakte (Woche 12-14)",
        ["Aufgabe", "PT", "Deliverable"],
        [
            ["CRM Grundkonfiguration", "2", "Leads & Opportunities"],
            ["Komplette Kundenhistorie", "3", "360° Kundensicht"],
            ["Aktivitäten & Aufgaben", "2", "Follow-ups"],
            ["Wiedervorlagen", "1", "Keine Leads vergessen"],
        ], ACCENT)
    
    # Lager & Buchhaltung (SPÄTER)
    add_section_slide(prs, "Phase 5-6: Lager & Finanzen", "🔧")
    
    add_table_slide(prs, "Lager & Bestand (Woche 15-17)",
        ["Aufgabe", "PT", "Zusatz"],
        [
            ["Mehrere Lagerorte", "2", "Eigenleistung möglich"],
            ["Bestandsführung & Umlagerung", "2", ""],
            ["Scanner-Integration", "3", "Bei mobilen Scannern"],
            ["Etikettendruck", "2", "Modul ~500€"],
            ["Meldebestand", "1", "Auto-Bestellung"],
        ], PRIMARY)
    
    add_table_slide(prs, "Buchhaltung & DATEV (Woche 18-19)",
        ["Aufgabe", "PT", "Zusatz"],
        [
            ["DATEV-Schnittstelle", "2", "Modul ~500€"],
            ["X-Rechnung / ZUGFeRD", "1", "inkl. OCR"],
            ["Eingangsrechnungsprüfung", "4", "Automatisiert"],
            ["Pflichtfelder & Validierung", "3", ""],
            ["Kontenplan & Migration", "2", ""],
        ], PRIMARY)
    
    # Komplexe Phasen
    add_section_slide(prs, "Phase 7-8: Service & Schnittstellen", "🔗")
    
    add_table_slide(prs, "Monteur-App & Service (Woche 20-21)",
        ["Aufgabe", "PT", "Deliverable"],
        [
            ["Wartungsverträge", "1", "Verträge digital"],
            ["Mobile App Monteure", "4", "Monteure arbeiten mobil"],
            ["Mängeldoku mit Foto", "2", "Fotos am Projekt"],
            ["Digitale Unterschrift", "1", "Kunde signiert digital"],
        ], DARK)
    
    add_table_slide(prs, "Schnittstellen & Go-Live (Woche 22-23)",
        ["Aufgabe", "PT", "Deliverable"],
        [
            ["S-Firm Schnittstelle", "2", "Überweisungen automatisiert"],
            ["GAEB-Import", "4", "Ausschreibungen → Angebote"],
            ["Schulung Key User", "3", "Team fit"],
            ["Go-Live & Hypercare", "3", "Produktiver Betrieb"],
        ], DARK)
    
    # Meilensteine
    add_table_slide(prs, "Meilensteine",
        ["Woche", "Datum", "Meilenstein"],
        [
            ["2", "21. Feb", "Analyse abgeschlossen"],
            ["5", "14. Mär", "Vertrieb produktiv"],
            ["8", "4. Apr", "Einkauf produktiv"],
            ["11", "25. Apr", "Projektcontrolling live"],
            ["14", "16. Mai", "CRM aktiv"],
            ["17", "6. Jun", "Lager läuft"],
            ["19", "20. Jun", "Buchhaltung komplett"],
            ["21", "4. Jul", "Monteure digital"],
            ["23", "20. Jul", "🚀 GO-LIVE"],
        ])
    
    # Nächste Schritte
    add_content_slide(prs, "Nächste Schritte", [
        "1. Kick-off Meeting: 9. Februar",
        "2. Key User pro Bereich benennen",
        "3. Server-Infrastruktur vorbereiten",
        "4. Datenexport aus Altsystem starten",
        "5. Hardware-Anforderungen klären",
    ])
    
    # Abschluss
    add_title_slide(prs, "Fragen?", "Becker Sicherheitstechnik\nOdoo ERP Projekt 2026")
    
    output_path = "/home/reisig/.openclaw/workspace/projects/becker-odoo/Projektplan_Becker_Odoo.pptx"
    prs.save(output_path)
    print(f"✅ Präsentation erstellt: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
